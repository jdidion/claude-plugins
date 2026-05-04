#!/Users/jodidion/.claude/scripts/.venv/bin/python
"""Convert a Typst slide deck to PPTX as image-per-slide.

Unlike MARP (which has structured markdown we can render into native
PPTX text/table/image shapes), a Typst deck is arbitrary code that
compiles to visual pages. The cleanest path to PPTX/Google-Slides is
therefore:

    typst compile --root <root> deck.typ tmp/page-{p}.png --format png --ppi <n>
    → read each PNG → one PPTX slide per page with that PNG sized to fill.

Trade-off: resulting slides are non-editable (each slide is one image).
That's unavoidable without a full Typst-AST-to-PPTX converter. Document
this in SKILL.md so users know what they're getting.
"""

from __future__ import annotations

import argparse
import re
import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

DEFAULT_PPI = 200


def find_typst_root(typ_path: Path) -> Path:
    """Walk up from typ_path looking for the project root.

    The --root flag is required whenever the .typ references files
    outside its own directory (e.g., `image("/diagrams/foo.png")` or
    `image("../diagrams/foo.png")`). We try, in order:

    1. A directory containing `typst.toml` (Typst's convention).
    2. A directory whose layout matches the absolute-style `image()`
       references in the source — e.g. if the source has
       `image("/diagrams/foo.png")` then the root is the first ancestor
       where `diagrams/foo.png` exists.
    3. Fallback: the typ's own directory.
    """
    d = typ_path.resolve().parent
    ancestors = [d, *d.parents]

    for candidate in ancestors:
        if (candidate / "typst.toml").exists():
            return candidate

    try:
        src = typ_path.read_text()
    except OSError:
        return d

    abs_refs = re.findall(r'image\(\s*"(/[^"]+)"', src)
    for candidate in ancestors:
        if all((candidate / ref.lstrip("/")).exists() for ref in abs_refs) and abs_refs:
            return candidate

    return d


def detect_aspect_ratio(typ_path: Path) -> tuple[float, float]:
    """Infer slide size in inches from the Typst source.

    Scans for `aspect-ratio: "16-9"` / `"4-3"` (Touying convention), or
    a `#set page(width: ..., height: ...)` directive. Falls back to
    16:9 at 13.333 × 7.5 inches.
    """
    src = typ_path.read_text()

    m = re.search(r'aspect-ratio:\s*"(\d+)[-:x](\d+)"', src)
    if m:
        w, h = float(m.group(1)), float(m.group(2))
        # Standard 16:9 PPTX slide is 13.333 x 7.5 in; scale to match ratio.
        slide_h = 7.5
        slide_w = slide_h * (w / h)
        return slide_w, slide_h

    m = re.search(
        r"#set\s+page\s*\(\s*width:\s*([\d.]+)\s*(cm|in|mm|pt)\s*,\s*height:\s*([\d.]+)\s*(cm|in|mm|pt)",
        src,
    )
    if m:
        def to_in(val: float, unit: str) -> float:
            return {
                "in": val,
                "cm": val / 2.54,
                "mm": val / 25.4,
                "pt": val / 72.0,
            }[unit]
        return to_in(float(m.group(1)), m.group(2)), to_in(float(m.group(3)), m.group(4))

    return 13.333, 7.5


def compile_typst_to_pngs(
    typ_path: Path, root: Path, out_dir: Path, ppi: int
) -> list[Path]:
    """Compile deck.typ to a sequence of page PNGs. Returns sorted paths."""
    pattern = out_dir / "page-{p}.png"
    cmd = [
        "typst", "compile",
        "--root", str(root),
        str(typ_path),
        str(pattern),
        "--format", "png",
        "--ppi", str(ppi),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        sys.stderr.write(result.stderr)
        raise SystemExit(f"typst compile failed: {result.returncode}")
    return sorted(out_dir.glob("page-*.png"), key=lambda p: int(re.search(r"page-(\d+)", p.name).group(1)))


def build_pptx(
    typ_path: Path,
    output_path: Path,
    root: Path | None = None,
    ppi: int = DEFAULT_PPI,
    slide_size_in: tuple[float, float] | None = None,
) -> None:
    typ_path = typ_path.resolve()
    if root is None:
        root = find_typst_root(typ_path)

    if slide_size_in is None:
        slide_size_in = detect_aspect_ratio(typ_path)
    slide_w_in, slide_h_in = slide_size_in

    with tempfile.TemporaryDirectory(prefix="typst-pptx-") as tmp:
        tmp_dir = Path(tmp)
        pngs = compile_typst_to_pngs(typ_path, root, tmp_dir, ppi)
        if not pngs:
            raise SystemExit("Typst compile produced no pages.")

        prs = Presentation()
        prs.slide_width = Emu(int(slide_w_in * 914400))
        prs.slide_height = Emu(int(slide_h_in * 914400))

        # Blank layout — last one is typically "Blank" in the default template.
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

        for png in pngs:
            slide = prs.slides.add_slide(blank_layout)
            # Size the image to fill the slide while preserving aspect.
            # Typst's aspect should match slide aspect, so letterboxing is minimal.
            with Image.open(png) as img:
                img_w, img_h = img.size
            img_aspect = img_w / img_h
            slide_aspect = slide_w_in / slide_h_in

            if abs(img_aspect - slide_aspect) < 0.01:
                # Aspects match — fill slide.
                left, top = Emu(0), Emu(0)
                width = prs.slide_width
                height = prs.slide_height
            elif img_aspect > slide_aspect:
                # Image wider — fit width, center vertically.
                width = prs.slide_width
                height = Emu(int(width / img_aspect))
                left = Emu(0)
                top = Emu(int((prs.slide_height - height) / 2))
            else:
                # Image taller — fit height, center horizontally.
                height = prs.slide_height
                width = Emu(int(height * img_aspect))
                top = Emu(0)
                left = Emu(int((prs.slide_width - width) / 2))

            slide.shapes.add_picture(str(png), left, top, width=width, height=height)

        prs.save(str(output_path))
        print(
            f"Saved {output_path} "
            f"({output_path.stat().st_size // 1024}KB, "
            f"{len(pngs)} slides, "
            f"{slide_w_in:.2f}×{slide_h_in:.2f} in @ {ppi} ppi)"
        )


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("input", type=Path, help="Path to deck.typ")
    ap.add_argument(
        "output", type=Path, nargs="?",
        help="Path to output .pptx (default: alongside input with .pptx extension)",
    )
    ap.add_argument(
        "--root", type=Path, default=None,
        help="Typst project root (default: walk up for typst.toml, else input's directory)",
    )
    ap.add_argument(
        "--ppi", type=int, default=DEFAULT_PPI,
        help=f"Render resolution in pixels per inch (default: {DEFAULT_PPI})",
    )
    ap.add_argument(
        "--aspect", type=str, default=None,
        help='Force slide aspect as "W:H" or "WxH" (default: inferred from .typ)',
    )
    args = ap.parse_args()

    if not args.input.exists():
        ap.error(f"input not found: {args.input}")
    output = args.output or args.input.with_suffix(".pptx")

    slide_size = None
    if args.aspect:
        m = re.match(r"(\d+(?:\.\d+)?)[:x](\d+(?:\.\d+)?)", args.aspect)
        if not m:
            ap.error(f"--aspect must look like W:H or WxH, got {args.aspect!r}")
        w, h = float(m.group(1)), float(m.group(2))
        slide_h = 7.5
        slide_size = (slide_h * (w / h), slide_h)

    build_pptx(args.input, output, root=args.root, ppi=args.ppi, slide_size_in=slide_size)
    return 0


if __name__ == "__main__":
    sys.exit(main())
