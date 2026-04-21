#!/Users/jodidion/.claude/scripts/.venv/bin/python
"""Inspect a PPTX template and guess which layouts play which roles.

Emits JSON with per-layout placeholder info and a best-guess role for each of
title / section / content / two_column. The /slides:template skill consumes
this output, prompts the user where confidence is low, then writes
~/.claude/scripts/slides_config.json.
"""

from __future__ import annotations

import argparse
import json
import sys
from dataclasses import asdict, dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER


# Placeholder roles we care about. Everything else is reported verbatim.
TITLE_TYPES = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE}
BODY_TYPES = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.PICTURE}
SUBTITLE_TYPES = {PP_PLACEHOLDER.SUBTITLE}


@dataclass
class PlaceholderInfo:
    idx: int
    ph_type: str
    name: str
    left_in: float
    top_in: float
    width_in: float
    height_in: float


@dataclass
class LayoutInfo:
    index: int
    name: str
    placeholders: list[PlaceholderInfo]
    role_guess: str
    confidence: str


def _ph_type_name(ph_type) -> str:
    try:
        return ph_type.name if ph_type else "NONE"
    except AttributeError:
        return str(ph_type)


def _emu_to_in(emu) -> float:
    return round(emu / 914400, 2) if emu is not None else 0.0


def _classify_layout(placeholders: list[PlaceholderInfo]) -> tuple[str, str]:
    """Guess a layout's role from its placeholders. Returns (role, confidence)."""
    decorative_types = {"SLIDE_NUMBER", "DATE", "FOOTER", "HEADER"}
    content_phs = [p for p in placeholders if p.ph_type not in decorative_types]
    titles = [p for p in content_phs if p.ph_type in {t.name for t in TITLE_TYPES}]
    bodies = [p for p in content_phs if p.ph_type in {t.name for t in BODY_TYPES}]
    subtitles = [p for p in content_phs if p.ph_type in {t.name for t in SUBTITLE_TYPES}]
    is_center_title = any(p.ph_type == PP_PLACEHOLDER.CENTER_TITLE.name for p in titles)

    # No content placeholders → blank
    if not content_phs:
        return ("blank", "high")

    # Title slide: centered title, optionally with subtitle. Decorative body
    # placeholders (author line etc.) are OK as long as title is centered.
    if is_center_title and not bodies:
        return ("title", "high")
    if is_center_title and subtitles and not bodies:
        return ("title", "high")
    if is_center_title:
        # Centered title with body-type placeholders — still a cover, medium
        return ("title", "medium")
    if titles and subtitles and not bodies and len(titles) == 1:
        return ("title", "medium")

    # Two-column: regular (non-centered) title + exactly two body placeholders,
    # side by side. Checked after title so covers don't steal this role.
    if titles and len(bodies) == 2:
        b1, b2 = sorted(bodies, key=lambda p: p.left_in)
        same_row = abs(b1.top_in - b2.top_in) < 0.5
        comparable_width = abs(b1.width_in - b2.width_in) < 1.0
        if same_row and comparable_width:
            return ("two_column", "high")
        return ("two_column", "medium")

    # Section header: title only, or title + small body text (divider)
    if titles and not bodies and not subtitles:
        return ("section", "medium")
    if titles and len(bodies) == 1 and bodies[0].height_in < 2.0:
        return ("section", "low")

    # Content: title + body placeholder(s)
    if titles and bodies:
        return ("content", "high" if len(bodies) == 1 else "medium")

    # Body-only (no title) → probably content with implicit title elsewhere
    if bodies and not titles:
        return ("content", "low")

    return ("unknown", "low")


def inspect(template_path: Path) -> dict:
    prs = Presentation(str(template_path))
    slide_w_in = _emu_to_in(prs.slide_width)
    slide_h_in = _emu_to_in(prs.slide_height)

    layouts: list[LayoutInfo] = []
    for i, layout in enumerate(prs.slide_layouts):
        phs: list[PlaceholderInfo] = []
        for ph in layout.placeholders:
            phs.append(
                PlaceholderInfo(
                    idx=ph.placeholder_format.idx,
                    ph_type=_ph_type_name(ph.placeholder_format.type),
                    name=ph.name,
                    left_in=_emu_to_in(ph.left),
                    top_in=_emu_to_in(ph.top),
                    width_in=_emu_to_in(ph.width),
                    height_in=_emu_to_in(ph.height),
                )
            )
        role, confidence = _classify_layout(phs)
        layouts.append(
            LayoutInfo(
                index=i,
                name=layout.name,
                placeholders=phs,
                role_guess=role,
                confidence=confidence,
            )
        )

    return {
        "template": str(template_path.resolve()),
        "slide_size_in": [slide_w_in, slide_h_in],
        "layout_count": len(layouts),
        "layouts": [asdict(L) for L in layouts],
        "role_picks": _pick_best_per_role(layouts),
    }


def _pick_best_per_role(layouts: list[LayoutInfo]) -> dict:
    """Pick the single best-matching layout per role, or None if no candidate."""
    priority = {"high": 3, "medium": 2, "low": 1}
    roles = ["title", "section", "content", "two_column"]
    subtitle_types = {t.name for t in SUBTITLE_TYPES}
    picks = {}
    for role in roles:
        candidates = [L for L in layouts if L.role_guess == role]
        if not candidates:
            picks[role] = None
            continue

        def score(L):
            s = (priority.get(L.confidence, 0), 0)
            # For title role, tiebreak toward layouts with a subtitle (real covers,
            # not quote/pull-quote slides that only have CENTER_TITLE).
            if role == "title":
                has_sub = any(p.ph_type in subtitle_types for p in L.placeholders)
                s = (s[0], 1 if has_sub else 0)
            return s

        candidates.sort(key=lambda L: (*score(L), -L.index), reverse=True)
        best = candidates[0]
        picks[role] = {
            "name": best.name,
            "index": best.index,
            "confidence": best.confidence,
        }
    return picks


def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("template", help="Path to .pptx template file")
    parser.add_argument("--pretty", action="store_true", help="Indent JSON output")
    args = parser.parse_args()

    path = Path(args.template).expanduser()
    if not path.exists():
        print(json.dumps({"error": f"file not found: {path}"}))
        sys.exit(1)
    if path.suffix.lower() != ".pptx":
        print(json.dumps({"error": f"not a .pptx file: {path}"}))
        sys.exit(1)

    result = inspect(path)
    print(json.dumps(result, indent=2 if args.pretty else None))


if __name__ == "__main__":
    main()
