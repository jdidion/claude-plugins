#!/Users/jodidion/.claude/scripts/.venv/bin/python
"""Convert MARP markdown slides to styled PPTX matching the teal/white theme."""

from __future__ import annotations

import json
import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# Theme colors
TEAL = RGBColor(0x00, 0x78, 0x8A)
DARK = RGBColor(0x2D, 0x34, 0x36)
GRAY = RGBColor(0x63, 0x6E, 0x72)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xDD, 0xDD, 0xDD)
CODE_BG = RGBColor(0xF5, 0xF6, 0xFA)

# Slide dimensions (16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_L = Inches(0.6)
MARGIN_R = Inches(0.6)
MARGIN_T = Inches(0.4)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R
HALF_W = Inches(5.8)
COL_GAP = Inches(0.4)

# HTML entity replacements
ENTITIES = {
    "&mdash;": "\u2014",
    "&rarr;": "\u2192",
    "&ge;": "\u2265",
    "&le;": "\u2264",
    "&lt;": "<",
    "&gt;": ">",
    "&amp;": "&",
}


def _find_layout(layout_map: dict, keywords: list[str], fallback=None):
    """Find a slide layout by trying keywords.

    Priority: exact match (all keywords) > substring match (all keywords) > fallback.
    """
    # Pass 1: exact match
    for kw in keywords:
        if kw in layout_map:
            return layout_map[kw]
    # Pass 2: case-insensitive substring
    names_lower = {name.lower(): name for name in layout_map}
    for kw in keywords:
        kw_lower = kw.lower()
        for name_lower, name in names_lower.items():
            if kw_lower in name_lower:
                return layout_map[name]
    return fallback


def _load_config() -> dict:
    """Load slides_config.json if it exists."""
    config_file = Path.home() / ".claude" / "scripts" / "slides_config.json"
    if config_file.exists():
        return json.loads(config_file.read_text())
    return {}


def _clear_placeholders(slide):
    """Remove inherited placeholder shapes so they don't show 'Click to add' text."""
    for ph in list(slide.placeholders):
        sp = ph._element
        sp.getparent().remove(sp)


def clean_text(text: str) -> str:
    """Replace HTML entities. Does NOT strip comments (handled per-slide)."""
    for ent, char in ENTITIES.items():
        text = text.replace(ent, char)
    return text.strip()


def add_styled_run(paragraph, text: str, bold=False, italic=False, size=Pt(18),
                    color=DARK, font_name="Calibri"):
    run = paragraph.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name


def add_rich_text(tf, text: str, size=Pt(18), color=DARK, alignment=PP_ALIGN.LEFT,
                   use_first_paragraph=False):
    """Add text with **bold** and `code` markdown rendered."""
    if use_first_paragraph and len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = Pt(4)

    # Split on **bold** and `code` patterns
    parts = re.split(r"(\*\*.*?\*\*|`[^`]+`)", text)
    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            add_styled_run(p, part[2:-2], bold=True, size=size, color=TEAL)
        elif part.startswith("`") and part.endswith("`"):
            add_styled_run(p, part[1:-1], size=Pt(size.pt - 2), color=DARK,
                          font_name="Courier New")
        elif part:
            add_styled_run(p, part, size=size, color=color)


def add_heading(tf, text: str, level: int, use_first_paragraph=False):
    """Add a heading to a text frame."""
    if use_first_paragraph and len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.space_after = Pt(8)
    text = text.lstrip("# ").strip()

    # Render bold markers within heading
    parts = re.split(r"(\*\*.*?\*\*)", text)
    match level:
        case 1:
            size, clr = Pt(32), TEAL
        case 2:
            size, clr = Pt(26), TEAL
        case _:
            size, clr = Pt(18), DARK

    for part in parts:
        if part.startswith("**") and part.endswith("**"):
            add_styled_run(p, part[2:-2], bold=True, size=size, color=clr)
        elif part:
            add_styled_run(p, part, bold=(level <= 2), size=size, color=clr)


def add_table(slide, rows_data: list[list[str]], left, top, width):
    """Add a styled table to the slide."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 0
    if n_rows == 0 or n_cols == 0:
        return top

    col_w = int(width / n_cols)
    row_h = Inches(0.4)
    table_h = row_h * n_rows

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, table_h)
    table = shape.table

    for i, row in enumerate(rows_data):
        for j, cell_text in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            display_text = cell_text.strip()

            if i == 0:  # Header row
                cell.fill.solid()
                cell.fill.fore_color.rgb = TEAL
                # Strip any bold markers from header
                plain = re.sub(r"\*\*(.+?)\*\*", r"\1", display_text)
                add_styled_run(p, plain, bold=True, size=Pt(14), color=WHITE)
            else:
                # Render **bold** segments in teal, rest in dark
                parts = re.split(r"(\*\*.*?\*\*)", display_text)
                for part in parts:
                    if part.startswith("**") and part.endswith("**"):
                        add_styled_run(p, part[2:-2], bold=True, size=Pt(14),
                                      color=TEAL)
                    elif part:
                        add_styled_run(p, part, size=Pt(14), color=DARK)

            # Cell margins
            cell.margin_left = Pt(6)
            cell.margin_right = Pt(6)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)

    return top + table_h + Inches(0.15)


def parse_table(lines: list[str]) -> list[list[str]]:
    """Parse markdown table lines into rows of cells."""
    rows = []
    for line in lines:
        line = line.strip()
        if not line.startswith("|"):
            continue
        cells = [c.strip() for c in line.split("|")[1:-1]]
        # Skip separator rows (---|---)
        if all(re.match(r"^[-:]+$", c) for c in cells):
            continue
        rows.append(cells)
    return rows


def parse_slide_content(text: str) -> list[dict]:
    """Parse slide text into content blocks."""
    blocks = []
    lines = text.split("\n")
    i = 0
    in_code = False
    code_lines = []

    while i < len(lines):
        line = lines[i]

        # Code blocks
        if line.strip().startswith("```"):
            if in_code:
                blocks.append({"type": "code", "text": "\n".join(code_lines)})
                code_lines = []
                in_code = False
            else:
                in_code = True
            i += 1
            continue
        if in_code:
            code_lines.append(line)
            i += 1
            continue

        # Skip HTML tags (divs, style blocks, script tags)
        stripped = line.strip()
        if stripped.startswith("<div") or stripped.startswith("</div"):
            i += 1
            continue
        if stripped.startswith("<style") or stripped.startswith("</style") or stripped.startswith("<script"):
            # Skip multi-line style blocks
            if stripped.startswith("<style") and ">" not in stripped:
                while i < len(lines) and "</style>" not in lines[i]:
                    i += 1
            i += 1
            continue

        # Headings
        if stripped.startswith("# "):
            blocks.append({"type": "h1", "text": stripped})
            i += 1
            continue
        if stripped.startswith("## "):
            blocks.append({"type": "h2", "text": stripped})
            i += 1
            continue
        if stripped.startswith("### "):
            blocks.append({"type": "h3", "text": stripped})
            i += 1
            continue

        # Tables
        if stripped.startswith("|"):
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i])
                i += 1
            blocks.append({"type": "table", "rows": parse_table(table_lines)})
            continue

        # Images
        img_match = re.match(r"!\[.*?\]\((.+?)\)", stripped)
        if img_match:
            blocks.append({"type": "image", "path": img_match.group(1)})
            i += 1
            continue

        # Bullet points
        if stripped.startswith("- ") or stripped.startswith("* "):
            blocks.append({"type": "bullet", "text": stripped[2:]})
            i += 1
            continue

        # Numbered lists
        num_match = re.match(r"^(\d+)\.\s+(.*)", stripped)
        if num_match:
            blocks.append({"type": "numbered", "num": num_match.group(1),
                          "text": num_match.group(2)})
            i += 1
            continue

        # Regular text (non-empty)
        if stripped:
            blocks.append({"type": "text", "text": stripped})

        i += 1

    return blocks


def detect_columns(text: str) -> tuple[str, str] | None:
    """If text has column layout, return (left_content, right_content)."""
    if '<div class="columns">' not in text:
        return None

    # Remove outer columns wrapper
    inner = re.sub(r'<div class="columns">\s*', "", text, count=1)

    # Split on the boundary between first </div><div> pair
    # Find the matching </div> for the first <div>
    parts = re.split(r"</div>\s*<div>", inner, maxsplit=1)
    if len(parts) != 2:
        return None

    left = re.sub(r"^\s*<div>\s*", "", parts[0])
    right = re.sub(r"\s*</div>\s*</div>\s*$", "", parts[1])
    return left, right


def _estimate_text_height(text: str, width_inches: float, font_pt: float) -> float:
    """Estimate height in inches for wrapped text."""
    # Approximate chars per line based on font size and column width
    chars_per_line = max(1, int(width_inches * 72 / (font_pt * 0.55)))
    # Strip markdown bold markers for length estimate
    plain = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    plain = re.sub(r"`(.+?)`", r"\1", plain)
    n_lines = max(1, -(-len(plain) // chars_per_line))  # ceil division
    line_height = font_pt / 72 * 1.4  # 1.4x line spacing
    return n_lines * line_height + 0.08  # small padding


def render_blocks(slide, blocks: list[dict], left, top, width, base_dir: Path):
    """Render content blocks into a region of the slide.

    Batches consecutive text/bullet/numbered items into single text frames
    so PowerPoint handles wrapping naturally.
    """
    y = top
    width_in = width / 914400  # EMU to inches

    # Collect consecutive text-like blocks into batches
    i = 0
    while i < len(blocks):
        block = blocks[i]

        match block["type"]:
            case "h1":
                tf = slide.shapes.add_textbox(left, y, width, Inches(0.45)).text_frame
                tf.word_wrap = True
                tf.auto_size = None  # fixed size, don't auto-expand
                add_heading(tf, block["text"], 1)
                y += Inches(0.7)
                i += 1

            case "h2":
                tf = slide.shapes.add_textbox(left, y, width, Inches(0.35)).text_frame
                tf.word_wrap = True
                tf.auto_size = None
                add_heading(tf, block["text"], 2)
                y += Inches(0.75)
                i += 1

            case "h3":
                tf = slide.shapes.add_textbox(left, y, width, Inches(0.35)).text_frame
                tf.word_wrap = True
                tf.auto_size = None
                add_heading(tf, block["text"], 3)
                y += Inches(0.45)
                i += 1

            case "table":
                y = add_table(slide, block["rows"], left, y, width)
                i += 1

            case "bullet" | "numbered" | "text":
                # Batch consecutive text-like blocks into one text frame
                batch = []
                total_h = 0.0
                while i < len(blocks) and blocks[i]["type"] in ("bullet", "numbered", "text"):
                    b = blocks[i]
                    text = clean_text(b.get("text", ""))
                    if not text:
                        i += 1
                        continue
                    if b["type"] == "bullet":
                        text = "\u2022 " + text
                    elif b["type"] == "numbered":
                        text = f"{b.get('num', '')}. " + text
                    batch.append(text)
                    total_h += _estimate_text_height(text, width_in, 18)
                    i += 1

                if batch:
                    tf = slide.shapes.add_textbox(
                        left, y, width, Inches(total_h)
                    ).text_frame
                    tf.word_wrap = True
                    for text in batch:
                        add_rich_text(tf, text, size=Pt(18))
                    y += Inches(total_h + 0.05)

            case "code":
                code_text = block["text"]
                n_lines = len(code_text.split("\n"))
                h = n_lines * 0.18 + 0.15
                box = slide.shapes.add_textbox(left, y, width, Inches(h))
                box.fill.solid()
                box.fill.fore_color.rgb = CODE_BG
                tf = box.text_frame
                tf.word_wrap = True
                p = tf.add_paragraph()
                add_styled_run(p, code_text, size=Pt(12), color=DARK,
                              font_name="Courier New")
                y += Inches(h + 0.1)
                i += 1

            case "image":
                img_path = base_dir / block["path"]
                if not img_path.exists():
                    img_path = base_dir.parent / "figures" / block["path"]
                if img_path.exists():
                    slide.shapes.add_picture(
                        str(img_path), left, y,
                        width=min(width, Inches(5)),
                    )
                    y += Inches(3.5)
                i += 1

            case _:
                i += 1

    return y


def build_pptx(
    md_path: str | Path,
    output_path: str | Path,
    template_path: str | Path | None = None,
):
    md_path = Path(md_path)
    output_path = Path(output_path)
    base_dir = md_path.parent

    with open(md_path) as f:
        content = f.read()

    # Strip MARP frontmatter
    parts = content.split("---\n", 2)
    if len(parts) >= 3 and "marp:" in parts[1]:
        content = parts[2]

    # Clean entities globally
    content = clean_text(content)

    # Split into slides
    slide_texts = re.split(r"\n---\n", content)

    # Load template or create blank presentation
    if template_path and Path(template_path).exists():
        prs = Presentation(str(template_path))
        # Delete all existing template slides (keep layouts/masters)
        for _ in range(len(prs.slides)):
            sldId = prs.slides._sldIdLst[0]
            rId = sldId.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
            if rId:
                prs.part.drop_rel(rId)
            prs.slides._sldIdLst.remove(sldId)
        print(f"Using template: {template_path} ({len(prs.slide_layouts)} layouts)")
    else:
        prs = Presentation()

    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Find layouts by name — try config overrides, then keyword search, then fallback
    layout_map = {sl.name: sl for sl in prs.slide_layouts}
    config = _load_config()
    layout_names = config.get("layout_names", {})

    def _resolve(role: str, keywords: list[str], fallback=None):
        """Resolve layout: config override > keyword search > fallback."""
        if role in layout_names and layout_names[role] in layout_map:
            return layout_map[layout_names[role]]
        return _find_layout(layout_map, keywords, fallback=fallback)

    last_layout = prs.slide_layouts[-1]
    blank_layout = _resolve("content", ["BLANK", "Blank", "Bullet", "Title_Subhead"], last_layout)
    title_layout = _resolve("title", ["TITLE", "1_Cover", "Cover"], blank_layout)
    section_layout = _resolve("section", ["SECTION_HEADER", "Divider"], title_layout)

    for slide_text in slide_texts:
        slide_text = slide_text.strip()
        if not slide_text:
            continue

        # Detect if this is a title/section slide (lead class + only headings)
        is_lead = "<!-- _class: lead -->" in slide_text
        slide_text_clean = re.sub(r"<!--.*?-->", "", slide_text, flags=re.DOTALL).strip()

        if is_lead:
            # Use TITLE layout for lead slides
            slide = prs.slides.add_slide(title_layout)
            blocks = parse_slide_content(slide_text_clean)
            # Fill title placeholder if available
            if title_layout != blank_layout and len(slide.placeholders) >= 2:
                title_ph = slide.placeholders[0]
                title_ph.text = ""
                h1_blocks = [b for b in blocks if b["type"] == "h1"]
                h2_blocks = [b for b in blocks if b["type"] == "h2"]
                text_blocks = [b for b in blocks if b["type"] == "text"]
                if h1_blocks:
                    add_heading(title_ph.text_frame, h1_blocks[0]["text"], 1,
                               use_first_paragraph=True)
                subtitle_ph = slide.placeholders[1]
                subtitle_ph.text = ""
                if h2_blocks:
                    add_heading(subtitle_ph.text_frame, h2_blocks[0]["text"], 2,
                               use_first_paragraph=True)
                if text_blocks:
                    add_rich_text(subtitle_ph.text_frame, clean_text(text_blocks[0]["text"]),
                                 size=Pt(18), color=GRAY, use_first_paragraph=not h2_blocks)
            else:
                # Fallback: manual text boxes
                y = Inches(1.5)
                for b in blocks:
                    if b["type"] == "h1":
                        tf = slide.shapes.add_textbox(MARGIN_L, y, CONTENT_W, Inches(0.7)).text_frame
                        tf.word_wrap = True
                        add_heading(tf, b["text"], 1, use_first_paragraph=True)
                        y += Inches(0.9)
                    elif b["type"] in ("h2", "text"):
                        tf = slide.shapes.add_textbox(MARGIN_L, y, CONTENT_W, Inches(0.5)).text_frame
                        tf.word_wrap = True
                        text = b["text"].lstrip("# ").strip() if b["type"] == "h2" else b["text"]
                        add_rich_text(tf, text, size=Pt(18), color=GRAY,
                                     use_first_paragraph=True)
                        y += Inches(0.6)
            continue

        # Check for two-column layout
        cols = detect_columns(slide_text_clean)

        if cols:
            left_text, right_text = cols
            slide = prs.slides.add_slide(blank_layout)
            _clear_placeholders(slide)

            # Find slide title (first # heading before or in columns)
            title_match = re.match(r"^(#\s+.+?)$", slide_text_clean, re.MULTILINE)
            title_y = MARGIN_T

            if title_match:
                title = title_match.group(1)
                tf = slide.shapes.add_textbox(MARGIN_L, title_y, CONTENT_W, Inches(0.45)).text_frame
                tf.word_wrap = True
                tf.auto_size = None
                add_heading(tf, title, 1, use_first_paragraph=True)
                title_y += Inches(0.85)
                # Remove title from column content
                left_text = re.sub(r"^#\s+.+?\n", "", left_text, count=1).strip()

            # Left column
            left_blocks = parse_slide_content(left_text)
            render_blocks(slide, left_blocks, MARGIN_L, title_y, HALF_W, base_dir)

            # Right column
            right_blocks = parse_slide_content(right_text)
            render_blocks(slide, right_blocks,
                         MARGIN_L + HALF_W + COL_GAP, title_y, HALF_W, base_dir)
        else:
            # Single column
            slide = prs.slides.add_slide(blank_layout)
            _clear_placeholders(slide)
            blocks = parse_slide_content(slide_text_clean)

            y = MARGIN_T
            # If first block is h1, render it as the slide title
            if blocks and blocks[0]["type"] == "h1":
                tf = slide.shapes.add_textbox(MARGIN_L, y, CONTENT_W, Inches(0.45)).text_frame
                tf.word_wrap = True
                tf.auto_size = None
                add_heading(tf, blocks[0]["text"], 1, use_first_paragraph=True)
                y += Inches(0.85)
                blocks = blocks[1:]

            render_blocks(slide, blocks, MARGIN_L, y, CONTENT_W, base_dir)

    prs.save(str(output_path))
    print(f"Saved {output_path} ({output_path.stat().st_size // 1024}KB)")


if __name__ == "__main__":
    import sys

    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    flags = {a.split("=")[0]: a.split("=", 1)[1] if "=" in a else True
             for a in sys.argv[1:] if a.startswith("--")}

    if len(args) < 1:
        print("Usage: marp_to_pptx.py <input.md> [output.pptx] [--template=path.pptx]")
        sys.exit(1)

    input_path = Path(args[0])
    output_path = Path(args[1]) if len(args) >= 2 else input_path.with_suffix(".pptx")

    # Template: explicit flag > config file > None
    template = flags.get("--template")
    if template is None:
        template = _load_config().get("default_template")

    build_pptx(input_path, output_path, template_path=template)
